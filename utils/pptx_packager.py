"""
PowerPoint Packager for RGBA Layers

Creates PowerPoint presentations with each layer on a separate slide.
"""

from pptx import Presentation
from pptx.util import Inches
from PIL import Image
from io import BytesIO
from typing import List, Dict, Any


def create_pptx(layers: List[Image.Image], params: Dict[str, Any]) -> bytes:
    """
    Create PowerPoint presentation with each layer on a separate slide

    Args:
        layers: List of RGBA PIL Images
        params: Parameters used for generation

    Returns:
        bytes: PPTX file as bytes
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(10)

    for i, layer in enumerate(layers):
        # Add blank slide
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Add layer title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9), Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = f"Layer {i}"

        # Convert PIL Image to bytes
        img_buffer = BytesIO()
        layer.save(img_buffer, format="PNG")
        img_buffer.seek(0)

        # Add image to slide
        slide.shapes.add_picture(
            img_buffer,
            Inches(1),
            Inches(1.5),
            width=Inches(8)
        )

    # Save to bytes
    output_buffer = BytesIO()
    prs.save(output_buffer)
    output_buffer.seek(0)
    return output_buffer.read()
